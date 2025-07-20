from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_slideshow_pptx():
    # Create a new presentation
    prs = Presentation()
    
    # Define slide layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    # Slide 1: Welcome with gradient background
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background shape
    background_shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill = background_shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)  # Deep blue
    fill.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)  # Burgundy
    background_shape.line.fill.background()
    
    title1 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    title1.text = "Welcome to the\nSBM & MVULANA Summit"
    title1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title1.text_frame.paragraphs[0].font.size = Pt(48)
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)  # Gold
    title1.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box1 = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(8), Inches(2)
    )
    content_box1.fill.solid()
    content_box1.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    content_box1.line.color.rgb = RGBColor(233, 185, 73)  # Gold
    content_box1.line.width = Pt(2)
    
    # Slide 2: SBM Logo with branding
    slide2 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape2 = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill2 = background_shape2.fill
    fill2.gradient()
    fill2.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill2.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape2.line.fill.background()
    
    title2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title2.text = "SBM"
    title2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title2.text_frame.paragraphs[0].font.size = Pt(48)
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title2.text_frame.paragraphs[0].font.bold = True
    
    # Add logo placeholder
    logo_placeholder = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2.5), Inches(4), Inches(3)
    )
    logo_placeholder.fill.solid()
    logo_placeholder.fill.fore_color.rgb = RGBColor(233, 185, 73)  # Gold
    logo_placeholder.line.color.rgb = RGBColor(233, 185, 73)
    logo_placeholder.line.width = Pt(3)
    
    # Slide 3: MVULANA with enhanced styling
    slide3 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape3 = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill3 = background_shape3.fill
    fill3.gradient()
    fill3.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill3.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape3.line.fill.background()
    
    title3 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title3.text = "MVULANA"
    title3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title3.text_frame.paragraphs[0].font.size = Pt(48)
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title3.text_frame.paragraphs[0].font.bold = True
    
    subtitle3 = slide3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    subtitle3.text = "Young Men's Leadership Initiative"
    subtitle3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle3.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle3.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box3 = slide3.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(2)
    )
    content_box3.fill.solid()
    content_box3.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box3.line.color.rgb = RGBColor(233, 185, 73)
    content_box3.line.width = Pt(2)
    
    content3 = slide3.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
    content3.text = "Empowering young men to become leaders in their communities."
    content3.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    content3.text_frame.paragraphs[0].font.size = Pt(18)
    content3.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content3.text_frame.paragraphs[0].font.bold = True
    
    # Slide 4: Bethune-Cookman University
    slide4 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape4 = slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill4 = background_shape4.fill
    fill4.gradient()
    fill4.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill4.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape4.line.fill.background()
    
    title4 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title4.text = "Bethune-Cookman University"
    title4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title4.text_frame.paragraphs[0].font.size = Pt(36)
    title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title4.text_frame.paragraphs[0].font.bold = True
    
    subtitle4 = slide4.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    subtitle4.text = "Center for Civic Engagement"
    subtitle4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle4.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle4.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle4.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box4 = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(2)
    )
    content_box4.fill.solid()
    content_box4.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box4.line.color.rgb = RGBColor(233, 185, 73)
    content_box4.line.width = Pt(2)
    
    content4 = slide4.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
    content4.text = "Proud host of the MVULANA Young Men's Summit"
    content4.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    content4.text_frame.paragraphs[0].font.size = Pt(18)
    content4.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content4.text_frame.paragraphs[0].font.bold = True
    
    # Slide 5: Workshop Overview with bullet points
    slide5 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape5 = slide5.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill5 = background_shape5.fill
    fill5.gradient()
    fill5.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill5.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape5.line.fill.background()
    
    title5 = slide5.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title5.text = "Workshop Overview"
    title5.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title5.text_frame.paragraphs[0].font.size = Pt(36)
    title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title5.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box5 = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(4.5)
    )
    content_box5.fill.solid()
    content_box5.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box5.line.color.rgb = RGBColor(233, 185, 73)
    content_box5.line.width = Pt(2)
    
    content5 = slide5.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3.5))
    content5.text = "Interactive sessions focused on personal and professional development:\n\n• Leadership Development & Confidence Building\n• Effective Communication Skills\n• Mental Health Awareness & Wellness\n• Career Planning & Goal Setting\n• Community Engagement & Service\n• Financial Literacy & Responsibility"
    content5.text_frame.paragraphs[0].font.size = Pt(18)
    content5.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content5.text_frame.paragraphs[0].font.bold = True
    
    # Slide 6: Event Schedule with table styling
    slide6 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape6 = slide6.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill6 = background_shape6.fill
    fill6.gradient()
    fill6.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill6.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape6.line.fill.background()
    
    title6 = slide6.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title6.text = "Event Schedule"
    title6.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title6.text_frame.paragraphs[0].font.size = Pt(36)
    title6.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title6.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box6 = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(5)
    )
    content_box6.fill.solid()
    content_box6.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box6.line.color.rgb = RGBColor(233, 185, 73)
    content_box6.line.width = Pt(2)
    
    content6 = slide6.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
    content6.text = "10:00-10:15 AM - Welcome/Prayer (Dr. McConner & Mr. McKinney/Bishop Triplett)\n10:15-10:30 AM - Ice Breaker (Mr. Jermaine McKinney)\n10:30-11:10 AM - Break Out Groups (1&2) Rotation (Group Facilitators)\n11:10-11:20 AM - Participants Break (Group Participants)\n11:30 AM-12:00 PM - Mayor Henry Presentation (Mayor of Daytona Beach Florida)\n12:00-12:30 PM - Lunch Break (Committee Members/Omega Psi Phi)\n12:30-1:10 PM - Break Out Groups (3&4) Rotation (Group Facilitators)\n12:30-1:10 PM - Gabriel Hannans - Parents Seminar: Gentle Parenting (Gabriel Hannans)\n1:10-1:40 PM - Black History Presentation (Dr. Headley White)\n1:40-2:20 PM - Break Out Groups (5&6) Rotation (Group Facilitators)\n2:20-2:50 PM - Power of the Tie Presentation (Mr. Mayner & Mr. Jermaine McKinney)\n2:50-4:00 PM - Participants Feedback/Presentation of Certificates & Closing Ceremony (Group Participants/Dr. McConner & Mr. McKinney)"
    content6.text_frame.paragraphs[0].font.size = Pt(12)
    content6.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content6.text_frame.paragraphs[0].font.bold = True
    
    # Slide 7: Kevin McCrary with speaker styling
    slide7 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape7 = slide7.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill7 = background_shape7.fill
    fill7.gradient()
    fill7.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill7.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape7.line.fill.background()
    
    title7 = slide7.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title7.text = "Group Facilitator"
    title7.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title7.text_frame.paragraphs[0].font.size = Pt(36)
    title7.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title7.text_frame.paragraphs[0].font.bold = True
    
    # Add speaker photo placeholder
    photo_placeholder7 = slide7.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(2), Inches(2)
    )
    photo_placeholder7.fill.solid()
    photo_placeholder7.fill.fore_color.rgb = RGBColor(233, 185, 73)
    photo_placeholder7.line.color.rgb = RGBColor(233, 185, 73)
    photo_placeholder7.line.width = Pt(3)
    
    # Add glassmorphic content box
    content_box7 = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2), Inches(5.5), Inches(4)
    )
    content_box7.fill.solid()
    content_box7.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box7.line.color.rgb = RGBColor(233, 185, 73)
    content_box7.line.width = Pt(2)
    
    content7 = slide7.shapes.add_textbox(Inches(4), Inches(2.5), Inches(4.5), Inches(3))
    content7.text = "Kevin McCrary\nDirector of Financial Aid Services, Bethune-Cookman University\n\nKevin McCrary is the Director of Financial Aid Services at Bethune-Cookman University in Daytona Beach, Florida. With nearly seven years of dedicated service to the institution, he has held various roles including Associate Director and Coordinator of Financial Aid Services.\n\nKevin is known for his commitment to student success and financial empowerment in higher education. His leadership and service have earned him recognition from professional organizations such as the National Association of Student Financial Aid Administrators (NASFAA)."
    content7.text_frame.paragraphs[0].font.size = Pt(14)
    content7.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content7.text_frame.paragraphs[0].font.bold = True
    
    # Slide 8: Gabriel Hannans
    slide8 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape8 = slide8.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill8 = background_shape8.fill
    fill8.gradient()
    fill8.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill8.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape8.line.fill.background()
    
    title8 = slide8.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title8.text = "Group Facilitator"
    title8.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title8.text_frame.paragraphs[0].font.size = Pt(36)
    title8.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title8.text_frame.paragraphs[0].font.bold = True
    
    # Add speaker photo placeholder
    photo_placeholder8 = slide8.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(2), Inches(2)
    )
    photo_placeholder8.fill.solid()
    photo_placeholder8.fill.fore_color.rgb = RGBColor(233, 185, 73)
    photo_placeholder8.line.color.rgb = RGBColor(233, 185, 73)
    photo_placeholder8.line.width = Pt(3)
    
    # Add glassmorphic content box
    content_box8 = slide8.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(2), Inches(5.5), Inches(4)
    )
    content_box8.fill.solid()
    content_box8.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box8.line.color.rgb = RGBColor(233, 185, 73)
    content_box8.line.width = Pt(2)
    
    content8 = slide8.shapes.add_textbox(Inches(4), Inches(2.5), Inches(4.5), Inches(3))
    content8.text = "Gabriel Hannans\nParent Coach & Educator, Creator of SHIFT Parenting\n\nGabriel Hannans is a nationally recognized parent coach, educator, and the creator of SHIFT Parenting—a practical, science-backed framework that helps overwhelmed caregivers raise emotionally healthy, self-regulated children without relying on punishment or shame.\n\nWith a background in Exceptional Student Education and years of experience as a registered behavior technician, Gabriel has supported countless neurodivergent children and families struggling with behavior, communication, and connection."
    content8.text_frame.paragraphs[0].font.size = Pt(14)
    content8.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content8.text_frame.paragraphs[0].font.bold = True
    
    # Slide 13: Derrick Henry - Speaker of the Day (special styling)
    slide13 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape13 = slide13.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill13 = background_shape13.fill
    fill13.gradient()
    fill13.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill13.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape13.line.fill.background()
    
    title13 = slide13.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title13.text = "Speaker of the Day"
    title13.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title13.text_frame.paragraphs[0].font.size = Pt(36)
    title13.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title13.text_frame.paragraphs[0].font.bold = True
    
    # Add special highlight box for Speaker of the Day
    highlight_box13 = slide13.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(9), Inches(0.8)
    )
    highlight_box13.fill.solid()
    highlight_box13.fill.fore_color.rgb = RGBColor(233, 185, 73)
    highlight_box13.line.color.rgb = RGBColor(233, 185, 73)
    highlight_box13.line.width = Pt(3)
    
    # Add speaker photo placeholder
    photo_placeholder13 = slide13.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1), Inches(3), Inches(2), Inches(2)
    )
    photo_placeholder13.fill.solid()
    photo_placeholder13.fill.fore_color.rgb = RGBColor(233, 185, 73)
    photo_placeholder13.line.color.rgb = RGBColor(233, 185, 73)
    photo_placeholder13.line.width = Pt(3)
    
    # Add glassmorphic content box
    content_box13 = slide13.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3), Inches(5.5), Inches(4)
    )
    content_box13.fill.solid()
    content_box13.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box13.line.color.rgb = RGBColor(233, 185, 73)
    content_box13.line.width = Pt(2)
    
    content13 = slide13.shapes.add_textbox(Inches(4), Inches(3.5), Inches(4.5), Inches(3))
    content13.text = "Derrick L. Henry\nMayor of Daytona Beach\n\nDerrick L. Henry, a Daytona Beach native, is serving his fourth term as mayor of the City of Daytona Beach, running unopposed in 2024. He is the city's 21st mayor and our community's second African American mayor.\n\nMayor Henry's platform \"Together we can do more\" focuses on improving quality of life, business development, youth mentoring, infrastructure upgrades, homelessness solutions, and affordable housing."
    content13.text_frame.paragraphs[0].font.size = Pt(14)
    content13.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content13.text_frame.paragraphs[0].font.bold = True
    
    # Slide 14: Registration with call-to-action styling
    slide14 = prs.slides.add_slide(blank_layout)
    
    # Add gradient background
    background_shape14 = slide14.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    fill14 = background_shape14.fill
    fill14.gradient()
    fill14.gradient_stops[0].color.rgb = RGBColor(26, 42, 71)
    fill14.gradient_stops[1].color.rgb = RGBColor(128, 0, 32)
    background_shape14.line.fill.background()
    
    title14 = slide14.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    title14.text = "Join Us"
    title14.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title14.text_frame.paragraphs[0].font.size = Pt(48)
    title14.text_frame.paragraphs[0].font.color.rgb = RGBColor(233, 185, 73)
    title14.text_frame.paragraphs[0].font.bold = True
    
    subtitle14 = slide14.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    subtitle14.text = "Registration Now Open"
    subtitle14.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle14.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle14.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle14.text_frame.paragraphs[0].font.bold = True
    
    # Add glassmorphic content box
    content_box14 = slide14.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(3)
    )
    content_box14.fill.solid()
    content_box14.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_box14.line.color.rgb = RGBColor(233, 185, 73)
    content_box14.line.width = Pt(2)
    
    content14 = slide14.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(2))
    content14.text = "Don't miss this opportunity to be part of a transformative experience!\n\nRegistration Deadline: July 15, 2025\nLimited Spaces Available\nFree Admission\n\nContact: [Your Contact Information]\nWebsite: [Your Website]"
    content14.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    content14.text_frame.paragraphs[0].font.size = Pt(18)
    content14.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    content14.text_frame.paragraphs[0].font.bold = True
    
    # Save the presentation
    prs.save('SBM-Mvulana-Summit-Styled-Content.pptx')
    print("PowerPoint presentation with styled content created successfully!")

if __name__ == "__main__":
    create_slideshow_pptx() 